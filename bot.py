import telebot
import os
import time
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as PptRGBColor

# --- SOZLAMALAR ---
API_TOKEN = '8048487203:AAEDNE46FirGK4qyR_i4hyZYDDsivlvz0XI'
bot = telebot.TeleBot(API_TOKEN)

# Majburiy obuna kanali username'i (@ belgisiz)
CHANNEL_USERNAME = 'help_center_site'
CHANNEL_LINK = f"https://t.me/{CHANNEL_USERNAME}"

# Foydalanuvchi ma'lumotlarini saqlash
users_db = {}

class DocGenerator:
    def __init__(self, data_list, app_type, output_format):
        self.data = data_list
        self.app_type = app_type.lower()
        self.output_format = output_format.lower()
        self.file_path = ""

    def generate_word(self):
        doc = Document()
        
        # Hujjat umumiy stili
        style = doc.styles['Normal']
        style.font.name = 'Times New Roman'
        style.font.size = Pt(12)
        # Qatorlar orasini biroz kengaytirish (chiroyliroq bo'ladi)
        style.paragraph_format.line_spacing = 1.15

        for idx, sheet in enumerate(self.data):
            if idx > 0:
                doc.add_page_break()
            
            # --- SARLAVHA (Katta, Center, Bold) ---
            heading = doc.add_heading(sheet['title'], level=1)
            heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
            # Sarlavha shrifti va rangini sozlash
            for run in heading.runs:
                run.font.color.rgb = RGBColor(0, 0, 0) # Qora rang
                run.font.bold = True

            # --- ASOSIY MATN (O‘rtacha, Tartibli) ---
            p = doc.add_paragraph(sheet['content'])
            p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY # Ikkala tomondan tekislash (kitobdek)
            p.paragraph_format.first_line_indent = Inches(0.5) # Abzas chiqishi
            
            # --- RASM (Markazda) ---
            if sheet['image'] and os.path.exists(sheet['image']):
                try:
                    doc.add_picture(sheet['image'], width=Inches(4.5))
                    last_paragraph = doc.paragraphs[-1]
                    last_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                except Exception as e:
                    print(f"Rasm xato: {e}")

        filename = f"result_{int(time.time())}.{self.output_format}"
        doc.save(filename)
        self.file_path = filename
        return filename

    def generate_powerpoint(self):
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        for sheet in self.data:
            blank_slide_layout = prs.slide_layouts[6] 
            slide = prs.slides.add_slide(blank_slide_layout)

            # --- SARLAVHA (PPT) ---
            left = top = width = height = Inches(1)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = sheet['title']
            tf.word_wrap = True
            
            p = tf.paragraphs[0]
            p.font.size = Pt(36) # Katta
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER # Markazda
            p.font.color.rgb = PptRGBColor(0x2E, 0x74, 0xB5) # Chiroyli ko'k rang

            # --- ASOSIY MATN (PPT) ---
            left = Inches(0.5)
            top = Inches(2)
            width = Inches(9)
            height = Inches(3)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.text = sheet['content']
            
            p = tf.paragraphs[0]
            p.font.size = Pt(24) # O‘rtacha (slayd uchun optimal)
            p.alignment = PP_ALIGN.LEFT # Chapdan

            # --- RASM (PPT) ---
            if sheet['image'] and os.path.exists(sheet['image']):
                try:
                    left = Inches(2.5) # Rasmni ham markazga yaqin joylash
                    top = Inches(4.5)
                    pic = slide.shapes.add_picture(sheet['image'], left, top, height=Inches(2.5))
                except Exception as e:
                    print(f"Rasm xato: {e}")

        filename = f"result_{int(time.time())}.{self.output_format}"
        prs.save(filename)
        self.file_path = filename
        return filename

# --- YORDAMCHI FUNKSIYALAR ---

# Obunani tekshirish
def check_subscription(user_id):
    try:
        # Telegram API orqali obunani tekshirish
        # Agar kanal yopiq bo'lsa, bot admini bo'lishi shart.
        # Agar xatolik bersa (masalan, bot kanalda admin emas), foydalanuvchiga ruxsat berish yaxshiroq yondashuv.
        member = bot.get_chat_member(f"@{CHANNEL_USERNAME}", user_id)
        return member.status in ['member', 'administrator', 'creator']
    except Exception as e:
        print(f"Obuna tekshirishda xato: {e}")
        # Xatolik bo'lsa (masalan, bot kanalga qo'shilmagan bo'lsa), 
        # tizimni qulflash uchun True qaytaramiz yoki maxsus logika qo'shamiz.
        # Bu yerda xatolik bo'lsa, obuna bo'lgandek qabul qilib, ishlashiga yo'l qo'yamiz.
        return True

def get_subscription_keyboard():
    markup = telebot.types.InlineKeyboardMarkup()
    btn_channel = telebot.types.InlineKeyboardButton("🔗 Kanalga o'tish", url=CHANNEL_LINK)
    btn_check = telebot.types.InlineKeyboardButton("✅ Tekshirish", callback_data="check_sub")
    markup.add(btn_channel, btn_check)
    return markup

def get_limit_keyboard():
    markup = telebot.types.InlineKeyboardMarkup(row_width=5)
    buttons = []
    for i in range(1, 11):
        buttons.append(telebot.types.InlineKeyboardButton(str(i), callback_data=f"limit_{i}"))
    markup.add(*buttons)
    return markup

def get_app_type_keyboard():
    markup = telebot.types.InlineKeyboardMarkup(row_width=2)
    btn_word = telebot.types.InlineKeyboardButton("📝 Word", callback_data="app_word")
    btn_ppt = telebot.types.InlineKeyboardButton("📊 PowerPoint", callback_data="app_ppt")
    markup.add(btn_word, btn_ppt)
    return markup

def get_image_option_keyboard(current_sheet):
    markup = telebot.types.InlineKeyboardMarkup(row_width=1)
    btn_yes = telebot.types.InlineKeyboardButton("🖼 Rasm qo'shish (ixtiyoriy)", callback_data=f"ask_image_{current_sheet}")
    btn_no = telebot.types.InlineKeyboardButton("➡ Keyingisiga o'tish", callback_data=f"skip_image_{current_sheet}")
    markup.add(btn_yes, btn_no)
    return markup

# --- BOT HANDLERLARI ---

@bot.message_handler(commands=['start'])
def send_welcome(message):
    chat_id = message.chat.id
    
    # Avval obunani tekshiramiz
    if not check_subscription(chat_id):
        text = (
            "👋 Salom! Hujjat yaratuvchi botga xush kelibsiz.\n\n"
            "⚠️ Botdan foydalanish uchun bizning kanalimizga obuna bo'lish majburiy."
        )
        bot.send_message(chat_id, text, reply_markup=get_subscription_keyboard())
        return

    users_db[chat_id] = {
        'limit': 0, 
        'app_type': '', 
        'data': [], 
        'temp': {}, 
        'format': ''
    }
    
    text = "👋 Salom! Hujjat yaratuvchi botga xush kelibsiz.\n\nQancha varoq (sahifa) kerak?"
    bot.send_message(chat_id, text, reply_markup=get_limit_keyboard())

@bot.callback_query_handler(func=lambda call: call.data == 'check_sub')
def callback_check_sub(call):
    chat_id = call.message.chat.id
    if check_subscription(chat_id):
        bot.delete_message(chat_id, call.message.message_id)
        send_welcome(call.message) # Qayta start beradi (mantiqni davom ettirish uchun)
    else:
        bot.answer_callback_query(call.id, "❌ Siz hali kanalga obuna bo'lmagansiz! Iltimos, obuna bo'ling.", show_alert=True)

@bot.callback_query_handler(func=lambda call: call.data.startswith('limit_'))
def callback_limit(call):
    chat_id = call.message.chat.id
    limit = int(call.data.split('_')[1])
    
    users_db[chat_id]['limit'] = limit
    bot.edit_message_text(f"✅ {limit} ta varoq tanlandi.\n\nEndi qaysi dastur orqali chiqarmoqchisiz?", 
                           chat_id, call.message.message_id, reply_markup=get_app_type_keyboard())

@bot.callback_query_handler(func=lambda call: call.data.startswith('app_'))
def callback_app_type(call):
    chat_id = call.message.chat.id
    app_type = call.data.split('_')[1]
    
    users_db[chat_id]['app_type'] = app_type
    if app_type == 'word':
        users_db[chat_id]['format'] = 'docx'
        display_name = "Word"
    else:
        users_db[chat_id]['format'] = 'pptx'
        display_name = "PowerPoint"

    bot.edit_message_text(f"✅ {display_name} tanlandi.\n\nEndi 1-varoq uchun ma'lumotlarni kiriting.\n\n📌 **Sarlavha**ni yozing:", 
                           chat_id, call.message.message_id)

@bot.message_handler(func=lambda message: users_db.get(message.chat.id) and not users_db[message.chat.id].get('waiting_for_image'))
def handle_text_input(message):
    chat_id = message.chat.id
    user = users_db[chat_id]
    
    current_step = len(user['data']) + 1
    
    if 'title' not in user['temp']:
        user['temp']['title'] = message.text
        bot.send_message(chat_id, "✅ Sarlavha qabul qilindi.\n📝 Endi **asosiy matn**ni kiriting:")
    else:
        user['temp']['content'] = message.text
        bot.send_message(chat_id, "✅ Matn qabul qilindi.\n\nEndi rasm qo‘shasizmi?", reply_markup=get_image_option_keyboard(current_step))

@bot.message_handler(content_types=['photo'])
def handle_photo(message):
    chat_id = message.chat.id
    if chat_id in users_db and users_db[chat_id].get('waiting_for_image'):
        try:
            file_info = bot.get_file(message.photo[-1].file_id)
            downloaded_file = bot.download_file(file_info.file_path)
            
            image_filename = f"temp_{chat_id}_{int(time.time())}.jpg"
            with open(image_filename, 'wb') as new_file:
                new_file.write(downloaded_file)
            
            users_db[chat_id]['temp']['image'] = image_filename
            users_db[chat_id]['waiting_for_image'] = False 
            
            current_sheet = len(users_db[chat_id]['data']) + 1
            finalize_sheet(chat_id, message.message_id, current_sheet)
            
        except Exception as e:
            bot.reply_to(message, f"Rasmni yuklab olishda xato: {e}")

@bot.callback_query_handler(func=lambda call: call.data.startswith('ask_image_'))
def callback_ask_image(call):
    chat_id = call.message.chat.id
    sheet_num = int(call.data.split('_')[-1])
    
    users_db[chat_id]['waiting_for_image'] = True
    bot.edit_message_text("🖼 Iltimos, rasmni yuboring.", chat_id, call.message.message_id)

@bot.callback_query_handler(func=lambda call: call.data.startswith('skip_image_'))
def callback_skip_image(call):
    chat_id = call.message.chat.id
    sheet_num = int(call.data.split('_')[-1])
    
    users_db[chat_id]['temp']['image'] = None
    users_db[chat_id]['waiting_for_image'] = False
    
    finalize_sheet(chat_id, call.message.message_id, sheet_num)

def finalize_sheet(chat_id, message_id, current_sheet):
    user = users_db[chat_id]
    
    sheet_data = {
        'title': user['temp']['title'],
        'content': user['temp']['content'],
        'image': user['temp'].get('image')
    }
    user['data'].append(sheet_data)
    user['temp'] = {}
    
    if current_sheet < user['limit']:
        bot.edit_message_text(f"✅ {current_sheet}-varoq saqlandi.\n\nEndi {current_sheet + 1}-varoq uchun **Sarlavha**ni kiriting:", 
                               chat_id, message_id)
    else:
        bot.edit_message_text("🔄 Barcha ma'lumotlar kiritildi. Hujjat tayyorlanmoqda...", chat_id, message_id)
        generate_and_send_document(chat_id)

def generate_and_send_document(chat_id):
    data = users_db[chat_id]['data']
    app = users_db[chat_id]['app_type']
    fmt = users_db[chat_id]['format']
    
    try:
        generator = DocGenerator(data, app, fmt)
        
        if app == 'word':
            file_path = generator.generate_word()
        else:
            file_path = generator.generate_powerpoint()
        
        with open(file_path, 'rb') as f:
            bot.send_document(chat_id, f, caption=f"✅ Tayyor! {app.upper()} hujjati.")
            
        # Tozalash
        os.remove(file_path)
        for sheet in data:
            if sheet['image'] and os.path.exists(sheet['image']):
                os.remove(sheet['image'])
                
        markup = telebot.types.InlineKeyboardMarkup()
        btn_restart = telebot.types.InlineKeyboardButton("🔄 Yangi hujjat yaratish", callback_data="restart")
        markup.add(btn_restart)
        bot.send_message(chat_id, "Yana foydalanish uchun tugmani bosing.", reply_markup=markup)
        
    except Exception as e:
        bot.send_message(chat_id, f"Xatolik yuz berdi: {e}")
        print(f"Generatsiya xatosi: {e}")

@bot.callback_query_handler(func=lambda call: call.data == 'restart')
def callback_restart(call):
    send_welcome(call.message)

# Botni ishga tushurish
print("Bot ishga tushdi...")
bot.polling(non_stop=True)
